from __future__ import annotations

from pathlib import Path
from typing import Iterable, Union
import logging

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover
    Presentation = None  # type: ignore

Logger = logging.getLogger(__name__)


class ImagesToPPTXError(RuntimeError):
    pass


def images_to_pptx(
    images: Iterable[Union[str, Path]],
    pptx_path: Union[str, Path],
    *,
    fit: str = "contain",  # "contain" | "cover"
) -> None:
    """
    Crée un PPTX avec une image par slide.

    - fit="contain" : image entière visible, ratio respecté (letterboxing possible)
    - fit="cover"   : image plein slide, ratio respecté, cropping possible
    """

    if Presentation is None:
        raise ImagesToPPTXError(
            "python-pptx non installé. Installez-le via `pip install python-pptx`"
        )

    pptx_path = Path(pptx_path)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    imgs = [Path(p) for p in images]
    if not imgs:
        raise ValueError("Aucune image fournie")

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank slide

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for img_path in imgs:
        if not img_path.exists():
            Logger.warning("Image introuvable, ignorée: %s", img_path)
            continue

        slide = prs.slides.add_slide(blank)

        try:
            pic = slide.shapes.add_picture(str(img_path), 0, 0)
        except Exception as exc:
            Logger.warning("Impossible d’ajouter l’image %s: %s", img_path, exc)
            continue

        img_w = pic.width
        img_h = pic.height

        # calcul du scaling
        scale_x = slide_w / img_w
        scale_y = slide_h / img_h

        if fit == "cover":
            scale = max(scale_x, scale_y)
        else:  # contain
            scale = min(scale_x, scale_y)

        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        pic.width = new_w
        pic.height = new_h

        # centrage
        pic.left = int((slide_w - new_w) / 2)
        pic.top = int((slide_h - new_h) / 2)

    if not prs.slides:
        raise ImagesToPPTXError("Aucune image valide n’a été ajoutée au PPTX")

    prs.save(str(pptx_path))


__all__ = ["images_to_pptx", "ImagesToPPTXError"]
