from pathlib import Path
import re
import random
from typing import Optional
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ==========================================================
# Style helpers
# ==========================================================

TITLE_FONT = "Calibri Light"
BODY_FONT = "Calibri"
TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(22)
BODY_SIZE = Pt(18)


def set_paragraph(p, text, size, bold=False):
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = size
    run.font.bold = bold


# ==========================================================
# Main converter
# ==========================================================

def _parse_attrs(s: str) -> dict:
    """Parse attribute block like {key=val key2=val2} into dict."""
    attrs = {}
    s = s.strip()
    if not s.startswith("{") or not s.endswith("}"):
        return attrs
    inner = s[1:-1].strip()
    for part in re.split(r"\s+", inner):
        if not part:
            continue
        if "=" in part:
            k, v = part.split("=", 1)
            attrs[k.strip()] = v.strip()
        else:
            attrs[part.strip()] = "true"
    return attrs


def _apply_style(prs: Presentation, style: str = "modern"):
    """Apply a simple global style seed (fonts, palette)."""
    # For now we only set a random accent color to use per-slide
    palettes = {
        "modern": [(0x2D, 0x6B, 0x9A), (0xF2, 0x7C, 0x5A)],
        "clean": [(0x1F, 0x1F, 0x1F), (0xFF, 0xFF, 0xFF)],
        "elegant": [(0x2E, 0x2A, 0x36), (0xD9, 0xC8, 0xB6)],
    }
    return palettes.get(style, palettes["modern"])


def markdown_to_pptx(md_path: str | Path, pptx_path: str | Path, *, max_slides: Optional[int] = None, style: str = "modern") -> None:
    md_path = Path(md_path)
    text = md_path.read_text(encoding="utf-8")

    prs = Presentation()

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    TITLE_ONLY = prs.slide_layouts[5]   # blank + title
    BLANK = prs.slide_layouts[6]

    current_title = None
    current_blocks = []
    current_images = []
    palettes = _apply_style(prs, style=style)
    accent_color = random.choice(palettes)

    def flush_slide():
        if not current_title:
            return

        slide = prs.slides.add_slide(BLANK)

        # apply subtle accent band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.4))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(*accent_color)
        band.line.fill.background()

        # ---------------- Title ----------------
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5),
            slide_width - Inches(1.4), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = current_title
        run.font.name = TITLE_FONT
        run.font.size = TITLE_SIZE
        run.font.bold = True

        # ---------------- Text content ----------------
        text_left = Inches(0.7)
        text_top = Inches(2.0)
        text_width = slide_width * 0.55
        text_height = slide_height - text_top - Inches(0.7)

        content_box = slide.shapes.add_textbox(
            text_left, text_top, text_width, text_height
        )
        tf = content_box.text_frame
        tf.clear()

        for block in current_blocks:
            if block["type"] == "subtitle":
                p = tf.add_paragraph()
                p.space_before = Pt(10)
                set_paragraph(p, block["text"], SUBTITLE_SIZE, bold=True)
            elif block["type"] == "bullet":
                p = tf.add_paragraph()
                p.level = block["level"]
                set_paragraph(p, block["text"], BODY_SIZE)

        # ---------------- Images ----------------
        if current_images:
            # current_images is a list of dicts: {path, pos, size}
            for img in current_images:
                img_path = Path(img["path"]) if isinstance(img, dict) else Path(img)
                if not img_path.exists():
                    continue
                pos = img.get("pos", "right") if isinstance(img, dict) else "right"
                size = img.get("size") if isinstance(img, dict) else None

                if pos == "full":
                    slide.shapes.add_picture(str(img_path), 0, 0, width=slide_width, height=slide_height)
                elif pos == "left":
                    img_left = text_left - Inches(0.05)
                    img_top = text_top
                    img_w = slide_width * 0.35
                    slide.shapes.add_picture(str(img_path), img_left, img_top, width=img_w)
                else:  # right/default
                    img_left = text_left + text_width + Inches(0.2)
                    img_top = text_top
                    img_w = slide_width - img_left - Inches(0.6)
                    slide.shapes.add_picture(str(img_path), img_left, img_top, width=img_w)

    # ======================================================
    # Parsing loop
    # ======================================================

    # support slide directives: <!-- slide: key=val, key2=val -->
    slide_directive_re = re.compile(r"<!--\s*slide:\s*(.*?)-->")

    for line in text.splitlines():
        line = line.rstrip()

        # Slide directive
        mdir = slide_directive_re.search(line)
        if mdir:
            # parse comma separated key=val
            raw = mdir.group(1)
            for part in re.split(r",\s*", raw):
                if "=" in part:
                    k, v = part.split("=", 1)
                    # allow background=image.jpg or style=clean
                    if k.strip() == "style":
                        palettes = _apply_style(prs, style=v.strip())
                        accent_color = random.choice(palettes)
            continue

        # New slide
        if line.startswith("# "):
            flush_slide()
            current_title = line[2:].strip()
            current_blocks = []
            current_images = []

        # Subtitle
        elif line.startswith("## "):
            current_blocks.append({
                "type": "subtitle",
                "text": line[3:].strip()
            })

        # Bullet (supports nesting via indentation)
        elif re.match(r"\s*-\s+", line):
            indent = len(line) - len(line.lstrip())
            level = indent // 2
            current_blocks.append({
                "type": "bullet",
                "text": line.strip()[2:].strip(),
                "level": min(level, 2)
            })

        # Image with optional attrs: ![alt](path){pos=right size=50%}
        elif m := re.match(r"!\[.*?\]\((.*?)\)(\{.*\})?", line):
            path = m.group(1)
            raw_attrs = m.group(2)
            attrs = _parse_attrs(raw_attrs) if raw_attrs else {}
            # slide target
            target = attrs.get("slide")
            entry = {"path": path}
            if "pos" in attrs:
                entry["pos"] = attrs["pos"]
            if "size" in attrs:
                entry["size"] = attrs["size"]
            if target:
                # if slide specified numerically, store as special directive in blocks
                try:
                    sidx = int(target)
                    current_blocks.append({"type": "attach_image_to_slide", "slide": sidx, "img": entry})
                except Exception:
                    current_images.append(entry)
            else:
                current_images.append(entry)

    flush_slide()

    pptx_path = Path(pptx_path)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


def convert(md_path: str | Path, pptx_path: str | Path, **kwargs) -> Path:
    markdown_to_pptx(md_path, pptx_path, **kwargs)
    return Path(pptx_path)


__all__ = ["markdown_to_pptx", "convert"]
