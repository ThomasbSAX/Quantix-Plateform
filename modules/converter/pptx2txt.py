from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from typing import List


def pptx_to_txt(
    pptx_path: str | Path,
    txt_path: str | Path,
    *,
    include_slide_numbers: bool = True,
    include_notes: bool = False,
    clean_whitespace: bool = True,
) -> Path:
    """
    Convertit un PPTX en texte structuré.

    Structure de sortie :

    === Slide 1 ===
    Titre
    - bullet
    - bullet

    === Slide 2 ===
    ...

    - Détecte titres / corps automatiquement
    - Ignore les shapes décoratifs
    """

    pptx_path = Path(pptx_path).resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)

    prs = Presentation(str(pptx_path))
    out_lines: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        # ─────────────────────────────
        # En-tête de slide
        # ─────────────────────────────
        if include_slide_numbers:
            out_lines.append(f"\n=== Slide {idx} ===\n")
        else:
            out_lines.append("\n=== Slide ===\n")

        title_text = None
        body_texts: List[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            # Détection titre (placeholder title)
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                title_text = text
            else:
                body_texts.append(text)

        if title_text:
            out_lines.append(title_text)
            out_lines.append("")

        for block in body_texts:
            lines = block.splitlines()
            for ln in lines:
                ln = ln.strip()
                if ln:
                    out_lines.append(f"- {ln}")

        # Notes du présentateur
        if include_notes and slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out_lines.append("\n[Notes]")
                out_lines.append(notes)

    text = "\n".join(out_lines)

    if clean_whitespace:
        while "\n\n\n" in text:
            text = text.replace("\n\n\n", "\n\n")
        text = text.strip() + "\n"

    txt_path = Path(txt_path).resolve()
    txt_path.parent.mkdir(parents=True, exist_ok=True)
    txt_path.write_text(text, encoding="utf-8")

    return txt_path


__all__ = ["pptx_to_txt"]


def convert(pptx_path: str | Path, txt_path: str | Path, **kwargs) -> Path:
    return pptx_to_txt(pptx_path, txt_path, **kwargs)


__all__.append("convert")
