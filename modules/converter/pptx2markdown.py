from __future__ import annotations

from pathlib import Path
from typing import List
from pptx import Presentation


def pptx_to_markdown_structured(
    pptx_path: str | Path,
    output_md: str | Path | None = None,
) -> Path:
    """
    Conversion PPTX → Markdown STRUCTURÉE.

    - 1 slide = 1 section Markdown
    - titre slide → #
    - sous-titre → ##
    - bullets → -
    """

    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)

    if output_md is None:
        output_md = pptx_path.parent / "result" / "markdown" / f"{pptx_path.stem}.md"
    output_md = Path(output_md)
    output_md.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    md_lines: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {slide_idx}"
        md_lines.append(f"# {title}\n")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue

            tf = shape.text_frame
            for para in tf.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # hiérarchie par niveau de bullet
                if para.level == 0:
                    md_lines.append(text)
                else:
                    md_lines.append("  " * para.level + f"- {text}")

        md_lines.append("")  # séparation slides

    output_md.write_text("\n".join(md_lines), encoding="utf-8")
    return output_md


__all__ = ["pptx_to_markdown_structured"]


def convert(pptx_path: str | Path, md_path: str | Path, **kwargs) -> Path:
    md_path = Path(md_path)
    produced = pptx_to_markdown_structured(pptx_path, output_md=md_path, **kwargs)
    return Path(produced)


__all__.append("convert")
